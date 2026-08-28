"""
CR-BioLM Thesis Presentation Builder
Professional academic presentation for expert demo.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color palette ─────────────────────────────────────────────────────────────
DARK_BG     = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
ACCENT      = RGBColor(0x22, 0xD3, 0xEE)   # cyan (matches habitat map)
ACCENT2     = RGBColor(0x38, 0xBD, 0xF8)   # lighter cyan
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY    = RGBColor(0x64, 0x74, 0x8B)
DARK_CARD   = RGBColor(0x1E, 0x29, 0x3B)   # card background
GREEN_OK    = RGBColor(0x34, 0xD3, 0x99)
AMBER       = RGBColor(0xFB, 0xBF, 0x24)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # blank
    return prs.slides.add_slide(layout)


def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_size=Pt(15), color=LIGHT_GRAY, spacing=Pt(6),
                  bold_first=False):
    """lines: list of str or (str, dict) tuples with overrides."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            text, opts = line
        else:
            text, opts = line, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        p.space_before = spacing
        run = p.add_run()
        run.text = text
        run.font.size = opts.get("size", font_size)
        run.font.bold = opts.get("bold", bold_first and i == 0)
        run.font.color.rgb = opts.get("color", color)
        run.font.italic = opts.get("italic", False)
    return txBox


def accent_bar(slide, top=Inches(0.08), height=Inches(0.055)):
    add_rect(slide, 0, top, SLIDE_W, height, fill_color=ACCENT)


def section_tag(slide, label, left=Inches(0.45), top=Inches(0.18)):
    add_rect(slide, left, top, Inches(2.4), Inches(0.32), fill_color=ACCENT)
    add_text(slide, label, left + Inches(0.12), top + Inches(0.01),
             Inches(2.2), Inches(0.3), font_size=Pt(11), bold=True,
             color=DARK_BG, align=PP_ALIGN.LEFT)


def slide_title(slide, title, sub=None, left=Inches(0.45), top=Inches(0.6)):
    add_text(slide, title, left, top, Inches(12), Inches(0.7),
             font_size=Pt(32), bold=True, color=WHITE)
    if sub:
        add_text(slide, sub, left, top + Inches(0.68), Inches(12), Inches(0.4),
                 font_size=Pt(16), color=ACCENT2)


def card(slide, left, top, width, height, title=None, lines=None,
         title_color=ACCENT, body_color=LIGHT_GRAY, body_size=Pt(14)):
    add_rect(slide, left, top, width, height, fill_color=DARK_CARD,
             line_color=ACCENT, line_width=Pt(0.75))
    cy = top + Inches(0.18)
    if title:
        add_text(slide, title, left + Inches(0.18), cy,
                 width - Inches(0.36), Inches(0.38),
                 font_size=Pt(14), bold=True, color=title_color)
        cy += Inches(0.38)
    if lines:
        add_multiline(slide, lines, left + Inches(0.18), cy,
                      width - Inches(0.36), height - (cy - top) - Inches(0.15),
                      font_size=body_size, color=body_color)


def bullet(text, indent=False, highlight=False):
    prefix = "    •  " if indent else "•  "
    return (prefix + text, {
        "color": ACCENT if highlight else LIGHT_GRAY,
        "bold": highlight,
    })


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

prs = new_prs()

# ── SLIDE 1: TÍTULO ───────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
accent_bar(s, top=0, height=Inches(0.12))
accent_bar(s, top=SLIDE_H - Inches(0.12), height=Inches(0.12))

add_rect(s, Inches(0), Inches(1.8), SLIDE_W, Inches(3.6), fill_color=DARK_CARD)

add_text(s, "CR-BioLM",
         Inches(0.6), Inches(2.0), Inches(12), Inches(1.1),
         font_size=Pt(72), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(s, "Modelado Multimodal de Distribución de Especies para Costa Rica",
         Inches(0.6), Inches(3.05), Inches(12), Inches(0.6),
         font_size=Pt(22), bold=False, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Random Forest · Explicabilidad SHAP · Modelos de Lenguaje con Visión",
         Inches(0.6), Inches(3.6), Inches(12), Inches(0.45),
         font_size=Pt(15), color=ACCENT2, align=PP_ALIGN.CENTER, italic=True)

add_text(s, "Tesis de Maestría — Demo del Marco de Evaluación  |  2026",
         Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
         font_size=Pt(13), color=MID_GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 2: PANORAMA GENERAL ─────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
accent_bar(s)
section_tag(s, "PANORAMA GENERAL")
slide_title(s, "¿Qué es CR-BioLM?",
            "Una tubería que pregunta: ¿puede un modelo de lenguaje con visión razonar sobre el hábitat de plantas?")

# 5 pilares
cols = [Inches(0.45), Inches(3.05), Inches(5.65), Inches(8.25), Inches(10.85)]
icons = ["📄", "🗺️", "🌿", "🤖", "📊"]
labels = ["Manual de\nPlantas CR", "Mapas de\nHábitat", "Catálogo de\n100 Especies", "Tubería\nLLM por Tiers", "Marco de\nEvaluación"]
descs = ["Extracción de\nficha técnica\nvía OCR + LLM",
         "43 polígonos\nbiogeográficos\nHammel + DEM",
         "4 volúmenes · mín.\n150 registros GBIF\npor especie",
         "T0 → T1 → T3\ntiers de contexto\ncreciente",
         "Juez ensamble\nrúbrica G-Eval\nCohen's κ"]

for i in range(5):
    add_rect(s, cols[i], Inches(2.4), Inches(2.45), Inches(3.9), fill_color=DARK_CARD,
             line_color=ACCENT, line_width=Pt(0.75))
    add_text(s, icons[i], cols[i], Inches(2.55), Inches(2.45), Inches(0.6),
             font_size=Pt(28), align=PP_ALIGN.CENTER)
    add_text(s, labels[i], cols[i], Inches(3.05), Inches(2.45), Inches(0.55),
             font_size=Pt(14), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, descs[i], cols[i], Inches(3.62), Inches(2.45), Inches(0.55),
             font_size=Pt(12), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "Pregunta de Investigación: ¿El contexto multimodal (mapas de hábitat + métricas climáticas) mejora la precisión del LLM en tareas de distribución de especies?",
         Inches(0.45), Inches(6.75), Inches(12.4), Inches(0.45),
         font_size=Pt(13), color=AMBER, italic=True, align=PP_ALIGN.CENTER)


# ── SLIDE 3: EXTRACCIÓN DE FICHA TÉCNICA ──────────────────────────────────────
s = blank_slide(prs)
bg(s)
accent_bar(s)
section_tag(s, "01 — EXTRACCIÓN DE DATOS")
slide_title(s, "Extracción de Datos de Hábitat del Manual de Plantas",
            "Extracción de texto OCR con PyMuPDF y estrategia de 3 niveles de respaldo")

# Izquierda: flujo del proceso
add_text(s, "Tubería de Extracción", Inches(0.45), Inches(1.75), Inches(5.5), Inches(0.35),
         font_size=Pt(14), bold=True, color=ACCENT)

steps = [
    ("1  Regex Exacto", "Localiza especie al inicio de línea · puntúa oraciones de hábitat por prioridad de palabras clave", GREEN_OK),
    ("2  Regex Difuso", "Maneja texto deformado por OCR · compara primeros 5 caracteres con tolerancia de comodín", AMBER),
    ("3  Respaldo LLM", "Llama 3.2 3B vía Groq (gratuito) · máx. 200 tokens · respeta límite de 8s por solicitud", ACCENT2),
]
for i, (title, desc, col) in enumerate(steps):
    top = Inches(2.15) + i * Inches(1.45)
    add_rect(s, Inches(0.45), top, Inches(5.5), Inches(1.3), fill_color=DARK_CARD,
             line_color=col, line_width=Pt(1.2))
    add_text(s, title, Inches(0.65), top + Inches(0.12), Inches(5.1), Inches(0.35),
             font_size=Pt(14), bold=True, color=col)
    add_text(s, desc, Inches(0.65), top + Inches(0.48), Inches(5.1), Inches(0.7),
             font_size=Pt(12), color=LIGHT_GRAY)

# Derecha: campos extraídos
add_text(s, "Campos Extraídos por Especie", Inches(6.4), Inches(1.75), Inches(6.4), Inches(0.35),
         font_size=Pt(14), bold=True, color=ACCENT)

fields = [
    ("habitat_raw", "Texto original extraído del PDF del Manual"),
    ("habitat_type", "Clasificación limpia (antes de la primera coma)"),
    ("elevation_min_m / max_m", "Extraído de notación '500–1500 m'"),
    ("elev_outlier_min/max", "Registros raros: '(0–)500–1850 (–2500) m'"),
    ("geographic_notes", "Referencias a vertiente, cordillera y cantón"),
    ("extraction_method", "exacto · difuso · llm · fallido"),
]
for i, (f, d) in enumerate(fields):
    top = Inches(2.15) + i * Inches(0.75)
    add_rect(s, Inches(6.4), top, Inches(6.4), Inches(0.68), fill_color=DARK_CARD,
             line_color=MID_GRAY, line_width=Pt(0.5))
    add_text(s, f, Inches(6.6), top + Inches(0.06), Inches(2.8), Inches(0.3),
             font_size=Pt(12), bold=True, color=ACCENT2)
    add_text(s, d, Inches(9.5), top + Inches(0.06), Inches(3.1), Inches(0.55),
             font_size=Pt(11), color=LIGHT_GRAY)

add_text(s, "118 reglas secuenciales de limpieza OCR corrigen saltos de línea con guión, confusiones dígito/letra y errores de codificación de acentos antes de la extracción.",
         Inches(0.45), Inches(6.8), Inches(12.4), Inches(0.38),
         font_size=Pt(11), color=MID_GRAY, italic=True)


# ── SLIDE 4: MAPA DE HÁBITAT ──────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
accent_bar(s)
section_tag(s, "02 — MAPEO DE HÁBITAT")
slide_title(s, "Construcción del Mapa de Hábitat",
            "Unidades Fitogeográficas Hammel (2014) + máscara DEM de elevación + puntos GBIF")

# Izquierda: capas de entrada
add_text(s, "Tres Capas de Entrada", Inches(0.45), Inches(1.75), Inches(4.5), Inches(0.35),
         font_size=Pt(14), bold=True, color=ACCENT)

layers = [
    ("🗺️  Shapefile Hammel", "43 polígonos biogeográficos\n18 grupos geográficos · cordilleras,\nllanuras, valles, penínsulas", ACCENT),
    ("⛰️  Ráster DEM", "altitud_cr.tif · máscara de elevación\ndentro de unidades coincidentes · señala\nregistros atípicos (zonas rayadas)", ACCENT2),
    ("📍  Presencias GBIF", "Puntos de ocurrencia filtrados para CR\nsuperpuestos como puntos rojos para\nvalidación con datos reales", GREEN_OK),
]
for i, (title, desc, col) in enumerate(layers):
    top = Inches(2.15) + i * Inches(1.5)
    add_rect(s, Inches(0.45), top, Inches(4.5), Inches(1.35), fill_color=DARK_CARD,
             line_color=col, line_width=Pt(1.0))
    add_text(s, title, Inches(0.65), top + Inches(0.1), Inches(4.1), Inches(0.35),
             font_size=Pt(13), bold=True, color=col)
    add_text(s, desc, Inches(0.65), top + Inches(0.45), Inches(4.1), Inches(0.8),
             font_size=Pt(12), color=LIGHT_GRAY)

# Centro: leyenda
add_text(s, "Codificación de Colores del Mapa", Inches(5.3), Inches(1.75), Inches(3.5), Inches(0.35),
         font_size=Pt(14), bold=True, color=ACCENT)

legend = [
    (RGBColor(0x2E, 0x34, 0x40), "Sin coincidencia geográfica"),
    (RGBColor(0x6B, 0x8A, 0x99), "Coincidencia geográfica · fuera del rango altitudinal"),
    (RGBColor(0x22, 0xD3, 0xEE), "Hábitat óptimo  (geografía + elevación ✓)"),
    (RGBColor(0x22, 0xD3, 0xEE), "Rayado — banda altitudinal atípica"),
    (RGBColor(0xEF, 0x44, 0x44), "Presencia GBIF confirmada (solo CR)"),
]
for i, (col, label) in enumerate(legend):
    top = Inches(2.2) + i * Inches(0.75)
    add_rect(s, Inches(5.3), top + Inches(0.1), Inches(0.4), Inches(0.35), fill_color=col,
             line_color=WHITE, line_width=Pt(0.3))
    add_text(s, label, Inches(5.85), top + Inches(0.06), Inches(3.0), Inches(0.4),
             font_size=Pt(12), color=LIGHT_GRAY)

# Derecha: lógica de emparejamiento
add_text(s, "Lógica de Emparejamiento", Inches(9.1), Inches(1.75), Inches(3.9), Inches(0.35),
         font_size=Pt(14), bold=True, color=ACCENT)

logic = [
    "Fase 1 — Recopilar códigos SUBUNIDAD\nde patrones de texto en geographic_notes\n(cordilleras, llanuras, valles...)",
    "Fase 2 — Filtro por vertiente\nSi solo se menciona Pacífico o Caribe,\nse elimina el lado opuesto de la cordillera",
    "Ejemplo:\n\"vert. Carib. Cord. de Guanacaste\"\n→ coincide solo con unidad 8.2 (no 8.1 Pac)",
]
for i, text in enumerate(logic):
    top = Inches(2.2) + i * Inches(1.45)
    add_rect(s, Inches(9.1), top, Inches(3.9), Inches(1.3), fill_color=DARK_CARD,
             line_color=MID_GRAY, line_width=Pt(0.5))
    add_text(s, text, Inches(9.3), top + Inches(0.12), Inches(3.5), Inches(1.1),
             font_size=Pt(12), color=LIGHT_GRAY)


# ── SLIDE 5: SELECCIÓN DE ESPECIES ────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
accent_bar(s)
section_tag(s, "03 — CATÁLOGO DE ESPECIES")
slide_title(s, "Catálogo Experimental de 100 Especies",
            "Selección estratificada en 4 volúmenes del Manual de Plantas de Costa Rica")

# Tarjetas de volumen
vols = [
    ("Vol. II", "Gimnospermas +\nMonocots A–M",
     "Agavaceae → Musaceae\nPinopsida, Cycadopsida,\nLiliopsida (parte)", "25 especies"),
    ("Vol. III", "Monocots O–Z",
     "Orchidaceae → Zingiberaceae\nLiliopsida (parte)\nIncluye Orquídeas, Bromelias", "25 especies"),
    ("Vol. VI", "Dicots H–P",
     "Haloragaceae → Phytolaccaceae\nMagnoliopsida (parte)\nIncluye Melastomataceae", "25 especies"),
    ("Vol. VIII", "Dicots S–Z",
     "Sabiaceae → Zygophyllaceae\nMagnoliopsida (parte)\nIncluye Urticaceae, Violaceae", "25 especies"),
]
for i, (vol, title, desc, count) in enumerate(vols):
    left = Inches(0.45) + i * Inches(3.22)
    add_rect(s, left, Inches(2.1), Inches(3.0), Inches(3.5), fill_color=DARK_CARD,
             line_color=ACCENT, line_width=Pt(0.75))
    add_text(s, vol, left + Inches(0.15), Inches(2.22), Inches(2.7), Inches(0.4),
             font_size=Pt(18), bold=True, color=ACCENT)
    add_text(s, title, left + Inches(0.15), Inches(2.62), Inches(2.7), Inches(0.45),
             font_size=Pt(13), bold=True, color=WHITE)
    add_text(s, desc, left + Inches(0.15), Inches(3.1), Inches(2.7), Inches(1.5),
             font_size=Pt(11), color=LIGHT_GRAY)
    add_rect(s, left + Inches(0.15), Inches(5.1), Inches(2.7), Inches(0.38), fill_color=ACCENT)
    add_text(s, count, left + Inches(0.15), Inches(5.12), Inches(2.7), Inches(0.35),
             font_size=Pt(14), bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# Caja de criterios de selección
add_rect(s, Inches(0.45), Inches(5.75), Inches(12.4), Inches(1.35), fill_color=DARK_CARD,
         line_color=AMBER, line_width=Pt(1.0))
add_text(s, "Criterios de Selección", Inches(0.65), Inches(5.87), Inches(4), Inches(0.35),
         font_size=Pt(13), bold=True, color=AMBER)
criteria = [
    "✓  Debe aparecer en el Manual de Plantas de Costa Rica (uno de los 4 volúmenes publicados)",
    "✓  Mínimo 150 registros de ocurrencia GBIF en Mesoamérica  (asegura datos suficientes para el RF)",
    "✓  Solo nombre binomial de especie  (excluye entradas solo de género y abreviaturas de autor)",
    "✓  25 especies muestreadas aleatoriamente por volumen  (semilla=42 para reproducibilidad)",
]
add_multiline(s, [(c, {"size": Pt(12), "color": LIGHT_GRAY}) for c in criteria],
              Inches(0.65), Inches(6.25), Inches(12.0), Inches(0.8))


# ── SLIDE 6: TUBERÍA LLM ──────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
accent_bar(s)
section_tag(s, "04 — TUBERÍA LLM")
slide_title(s, "Arquitectura de Tres Tiers de Contexto",
            "Inyección progresiva de información para medir la mejora del razonamiento del LLM")

# Cajas de tiers
tiers = [
    {
        "id": "T0", "name": "Línea Base de\nConocimiento Previo",
        "color": MID_GRAY,
        "input": "Solo nombre de la especie",
        "context": "Sin imágenes · sin datos\nSolo conocimiento de entrenamiento del LLM",
        "purpose": "Mide el conocimiento\nintríseco del LLM sobre\nla flora de Costa Rica",
    },
    {
        "id": "T1", "name": "Mapa de Distribución\nGBIF",
        "color": ACCENT2,
        "input": "1 imagen: mapa de ocurrencias\nGBIF en Mesoamérica",
        "context": "Regla estricta: razonar SOLO\ndesde la imagen · sin\nconocimiento previo permitido",
        "purpose": "Mide razonamiento\nespacial a partir de\ndatos de presencia",
    },
    {
        "id": "T3", "name": "Contexto Multimodal\nCompleto",
        "color": ACCENT,
        "input": "2 imágenes + 5 métricas",
        "context": "Mapa de hábitat Manual ·\nmapa predictivo RF ·\nAUC · elevación · SHAP",
        "purpose": "Mide razonamiento\nintegrado entre\nmúltiples fuentes",
    },
]

for i, t in enumerate(tiers):
    left = Inches(0.45) + i * Inches(4.3)
    col = t["color"]
    add_rect(s, left, Inches(1.85), Inches(4.05), Inches(4.8), fill_color=DARK_CARD,
             line_color=col, line_width=Pt(1.5))
    add_rect(s, left, Inches(1.85), Inches(4.05), Inches(0.65), fill_color=col)
    add_text(s, t["id"], left + Inches(0.15), Inches(1.88), Inches(0.7), Inches(0.58),
             font_size=Pt(28), bold=True, color=DARK_BG)
    add_text(s, t["name"], left + Inches(0.85), Inches(1.9), Inches(3.1), Inches(0.58),
             font_size=Pt(13), bold=True, color=DARK_BG)
    for j, (label, val) in enumerate([("ENTRADA", t["input"]), ("CONTEXTO", t["context"]), ("MIDE", t["purpose"])]):
        ty = Inches(2.6) + j * Inches(1.3)
        add_text(s, label, left + Inches(0.18), ty, Inches(3.6), Inches(0.28),
                 font_size=Pt(10), bold=True, color=col)
        add_text(s, val, left + Inches(0.18), ty + Inches(0.28), Inches(3.6), Inches(0.95),
                 font_size=Pt(12), color=LIGHT_GRAY)

# Flechas
for i in range(2):
    lx = Inches(4.5) + i * Inches(4.3)
    add_text(s, "→", lx, Inches(3.9), Inches(0.3), Inches(0.5),
             font_size=Pt(30), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Abajo: modelos
add_rect(s, Inches(0.45), Inches(6.85), Inches(12.4), Inches(0.42), fill_color=DARK_CARD,
         line_color=MID_GRAY, line_width=Pt(0.5))
add_text(s, "Modelos Generadores:  GPT-4o (OpenAI)   ·   Claude Sonnet 4.5 (Anthropic)   ·   Ambos llamados vía OpenRouter para cada tier × especie",
         Inches(0.65), Inches(6.88), Inches(12.0), Inches(0.36),
         font_size=Pt(12), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "T2 archivado — la inyección directa del texto de la ficha del Manual causó fuga del ground truth en los prompts de generación",
         Inches(0.45), Inches(6.5), Inches(12.4), Inches(0.3),
         font_size=Pt(10), color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ── SLIDE 7: MARCO DE EVALUACIÓN ─────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
accent_bar(s)
section_tag(s, "05 — MARCO DE EVALUACIÓN")
slide_title(s, "Juez LLM en Ensamble  +  Validación por Expertos Humanos",
            "Rúbrica G-Eval (Liu et al., 2023) con Evaluación de Ruta Condicional")

# Izquierda: métricas
add_text(s, "Métricas de Evaluación", Inches(0.45), Inches(1.75), Inches(6.0), Inches(0.35),
         font_size=Pt(14), bold=True, color=ACCENT)

metrics = [
    ("M1", "Precisión Geográfica",     "0–3", "TODOS", "¿Las zonas mencionadas (vertiente, cordillera) son consistentes con el Manual?"),
    ("M3", "Relevancia de Respuesta",  "0–3", "TODOS", "¿La respuesta aborda directamente la pregunta del usuario?"),
    ("M5", "Profundidad Analítica",    "0–3", "TODOS", "Calidad del razonamiento · integra fuentes · cobertura epistémica"),
    ("M2", "Uso del Contexto Altitudinal", "0–2", "T3", "Uso responsable de datos de elevación · reconcilia GBIF vs Manual"),
    ("M4", "Uso de Variable Climática","0–2", "T3",  "Variable SHAP correcta identificada · mecanismo ecológico explicado"),
]
for i, (mid, name, scale, tier, desc) in enumerate(metrics):
    top = Inches(2.15) + i * Inches(0.88)
    tier_col = ACCENT if tier == "TODOS" else AMBER
    add_rect(s, Inches(0.45), top, Inches(6.3), Inches(0.8), fill_color=DARK_CARD,
             line_color=MID_GRAY, line_width=Pt(0.4))
    add_text(s, mid, Inches(0.6), top + Inches(0.1), Inches(0.5), Inches(0.55),
             font_size=Pt(16), bold=True, color=ACCENT)
    add_text(s, name, Inches(1.15), top + Inches(0.1), Inches(2.5), Inches(0.3),
             font_size=Pt(12), bold=True, color=WHITE)
    add_text(s, scale, Inches(3.7), top + Inches(0.1), Inches(0.6), Inches(0.3),
             font_size=Pt(11), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(4.35), top + Inches(0.12), Inches(0.85), Inches(0.26), fill_color=tier_col)
    add_text(s, tier, Inches(4.35), top + Inches(0.12), Inches(0.85), Inches(0.26),
             font_size=Pt(10), bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(s, desc, Inches(1.15), top + Inches(0.42), Inches(5.5), Inches(0.32),
             font_size=Pt(10), color=LIGHT_GRAY)

# Fórmula de score
add_rect(s, Inches(0.45), Inches(6.55), Inches(6.3), Inches(0.65), fill_color=DARK_CARD,
         line_color=ACCENT, line_width=Pt(0.75))
add_text(s, "Compuesto:  T0/T1 = (M1+M3+M5) / 9    ·    T3 completo = (M1+M3+M5+M2+M4) / 13",
         Inches(0.65), Inches(6.65), Inches(6.0), Inches(0.42),
         font_size=Pt(12), bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# Derecha: ensamble de jueces
add_text(s, "Arquitectura del Ensamble de Jueces", Inches(7.1), Inches(1.75), Inches(5.9), Inches(0.35),
         font_size=Pt(14), bold=True, color=ACCENT)

judges = [
    ("Juez A", "Gemini 2.0 Flash", "OpenRouter", GREEN_OK),
    ("Juez B", "GPT-4o-mini",      "OpenRouter", ACCENT2),
    ("Juez C", "Llama 3.3 70B",    "Groq (gratuito) · desempate", AMBER),
]
for i, (jid, model, provider, col) in enumerate(judges):
    top = Inches(2.15) + i * Inches(0.85)
    add_rect(s, Inches(7.1), top, Inches(5.8), Inches(0.75), fill_color=DARK_CARD,
             line_color=col, line_width=Pt(1.0))
    add_text(s, jid, Inches(7.28), top + Inches(0.08), Inches(1.0), Inches(0.35),
             font_size=Pt(13), bold=True, color=col)
    add_text(s, model, Inches(8.35), top + Inches(0.08), Inches(2.5), Inches(0.35),
             font_size=Pt(13), color=WHITE)
    add_text(s, provider, Inches(7.28), top + Inches(0.43), Inches(4.5), Inches(0.25),
             font_size=Pt(10), color=MID_GRAY)

# Reglas
rules = [
    ("Guarda contra auto-favorecimiento", "Si familia del generador == familia del juez → se reemplaza ese juez por Llama C"),
    ("Bandera de desacuerdo", "|Juez A − Juez B| ≥ 2 en cualquier métrica → invocar C como desempate + marcar para revisión"),
    ("Validación humana", "Puntuaciones Likert de expertos (1–5) recolectadas vía formulario web · Cohen's κ ponderado lineal"),
]
for i, (title, desc) in enumerate(rules):
    top = Inches(4.8) + i * Inches(0.65)
    add_text(s, f"▸  {title}:", Inches(7.1), top, Inches(2.5), Inches(0.28),
             font_size=Pt(11), bold=True, color=ACCENT2)
    add_text(s, desc, Inches(9.7), top, Inches(3.2), Inches(0.55),
             font_size=Pt(11), color=LIGHT_GRAY)

add_rect(s, Inches(7.1), Inches(6.55), Inches(5.8), Inches(0.65), fill_color=DARK_CARD,
         line_color=AMBER, line_width=Pt(0.75))
add_text(s, "Hipótesis: Score T3 significativamente > T0/T1  →  el contexto multimodal mejora la precisión del LLM\nResultado preliminar (15 especies): T3 ≈ 0.84  vs  T1 ≈ 0.52  vs  T0 ≈ 0.61",
         Inches(7.3), Inches(6.6), Inches(5.4), Inches(0.56),
         font_size=Pt(11), color=AMBER, align=PP_ALIGN.CENTER)


# ── SLIDE 8: CIERRE ───────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s)
accent_bar(s, top=0, height=Inches(0.12))
accent_bar(s, top=SLIDE_H - Inches(0.12), height=Inches(0.12))

add_rect(s, Inches(0), Inches(2.0), SLIDE_W, Inches(3.5), fill_color=DARK_CARD)

add_text(s, "Próximos Pasos", Inches(0.6), Inches(2.15), Inches(12), Inches(0.5),
         font_size=Pt(24), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

next_steps = [
    "Corrida completa de 100 especies  ·  conjunto de datos experimental completo",
    "Evaluación por expertos humanos vía formulario web  (cr-biolm-production.up.railway.app)",
    "Análisis de concordancia Cohen's κ  ·  ANOVA entre tiers  ·  capítulo estadístico de la tesis",
    "Informe final y entrega de la tesis",
]
for i, step in enumerate(next_steps):
    add_text(s, f"{'①②③④'[i]}  {step}",
             Inches(1.5), Inches(2.85) + i * Inches(0.6), Inches(10), Inches(0.5),
             font_size=Pt(15), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "CR-BioLM  ·  Modelo de Lenguaje Biogeográfico para Costa Rica",
         Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
         font_size=Pt(13), color=MID_GRAY, align=PP_ALIGN.CENTER)


# ── GUARDAR ───────────────────────────────────────────────────────────────────
out = "g:/My Drive/Maestria/Tesis/CR-BioLM_Presentacion.pptx"
prs.save(out)
print(f"Guardado: {out}")
